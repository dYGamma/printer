import os
import sys
import logging
import asyncio
from PIL import Image, ImageDraw
import io
import json

logger = logging.getLogger(__name__)

# Try to import PDF library for preview
try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None
    logger.warning("pdf2image not available. PDF previews will be simulated.")


# Print settings defaults
DEFAULT_PRINT_SETTINGS = {
    "orientation": "auto",  # auto, portrait, landscape
    "scale": "fit",  # fit, fill, stretch, custom
    "custom_scale": 100,  # percentage for custom scale
    "center": True,
    "color_mode": "color"  # color, grayscale, black_white
}


def get_default_settings():
    """Returns a copy of default print settings."""
    return DEFAULT_PRINT_SETTINGS.copy()


async def create_preview(filepath: str, max_width: int = 800, max_height: int = 800) -> bytes:
    """
    Creates a preview image of a document for display in Telegram.
    Supports PDF, images (JPG, PNG), and Office documents.
    Returns image bytes that can be sent to Telegram.
    """
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext == '.pdf':
            return await _preview_pdf(filepath, max_width, max_height)
        elif ext in ['.jpg', '.jpeg', '.png']:
            return await _preview_image(filepath, max_width, max_height)
        elif ext in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt']:
            return await _preview_office_or_text(filepath, max_width, max_height)
        else:
            return _create_unsupported_preview(ext, max_width, max_height)
    except Exception as e:
        logger.error(f"Failed to create preview for {filepath}: {e}")
        return _create_error_preview(str(e), max_width, max_height)


async def _preview_pdf(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates preview from first page of PDF.
    """
    if convert_from_path is None:
        logger.warning("pdf2image not available, returning text preview for PDF")
        return _create_file_info_preview("PDF Document", filepath, max_width, max_height)
    
    try:
        loop = asyncio.get_event_loop()
        # Convert first page of PDF to image
        images = await loop.run_in_executor(
            None, 
            lambda: convert_from_path(filepath, first_page=1, last_page=1, dpi=150)
        )
        
        if images:
            img = images[0]
            return await _resize_and_encode_image(img, max_width, max_height)
    except Exception as e:
        logger.warning(f"Failed to convert PDF to preview: {e}")
        return _create_file_info_preview("PDF Document", filepath, max_width, max_height)


async def _preview_image(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates preview from image file.
    """
    try:
        loop = asyncio.get_event_loop()
        img = await loop.run_in_executor(None, Image.open, filepath)
        
        # Convert RGBA to RGB if needed
        if img.mode in ('RGBA', 'LA', 'P'):
            # Create white background
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        
        return await _resize_and_encode_image(img, max_width, max_height)
    except Exception as e:
        logger.error(f"Failed to preview image: {e}")
        return _create_error_preview(str(e), max_width, max_height)


async def _preview_office_or_text(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates text-based preview for Office documents and text files.
    """
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']:
            return await _preview_office(filepath, max_width, max_height)
        elif ext == '.txt':
            return await _preview_text(filepath, max_width, max_height)
    except Exception as e:
        logger.error(f"Failed to preview office/text: {e}")
    
    return _create_file_info_preview(ext.upper(), filepath, max_width, max_height)


async def _preview_office(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates preview for Office documents.
    """
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext in ['.docx', '.doc']:
            return await _preview_word(filepath, max_width, max_height)
        elif ext in ['.xlsx', '.xls']:
            return await _preview_excel(filepath, max_width, max_height)
        elif ext in ['.pptx', '.ppt']:
            return await _preview_powerpoint(filepath, max_width, max_height)
    except Exception as e:
        logger.warning(f"Could not preview office doc: {e}")
    
    return _create_file_info_preview(ext.upper(), filepath, max_width, max_height)


async def _preview_word(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates preview from Word document using python-docx.
    """
    try:
        from docx import Document
        loop = asyncio.get_event_loop()
        
        async def read_word():
            doc = Document(filepath)
            text_preview = "\n".join([p.text for p in doc.paragraphs[:10]])
            return text_preview
        
        text = await loop.run_in_executor(None, read_word)
        return _create_text_preview(text[:500], "Word Document", max_width, max_height)
    except ImportError:
        logger.warning("python-docx not available")
        return _create_file_info_preview("DOCX", filepath, max_width, max_height)


async def _preview_excel(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates preview from Excel document.
    """
    try:
        from openpyxl import load_workbook
        loop = asyncio.get_event_loop()
        
        async def read_excel():
            wb = load_workbook(filepath)
            ws = wb.active
            lines = []
            for row in ws.iter_rows(min_row=1, max_row=10, values_only=True):
                line = " | ".join(str(cell) if cell else "" for cell in row[:5])
                lines.append(line)
            return "\n".join(lines)
        
        text = await loop.run_in_executor(None, read_excel)
        return _create_text_preview(text[:500], "Excel Document", max_width, max_height)
    except ImportError:
        logger.warning("openpyxl not available")
        return _create_file_info_preview("XLSX", filepath, max_width, max_height)


async def _preview_powerpoint(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates text preview from PowerPoint.
    """
    try:
        from pptx import Presentation
        loop = asyncio.get_event_loop()
        
        async def read_pptx():
            prs = Presentation(filepath)
            texts = []
            for i, slide in enumerate(prs.slides[:3]):
                slide_text = f"Слайд {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                texts.append(slide_text)
            return "\n".join(texts)
        
        text = await loop.run_in_executor(None, read_pptx)
        return _create_text_preview(text[:500], "PowerPoint", max_width, max_height)
    except ImportError:
        logger.warning("python-pptx not available")
        return _create_file_info_preview("PPTX", filepath, max_width, max_height)


async def _preview_text(filepath: str, max_width: int, max_height: int) -> bytes:
    """
    Creates preview from text file.
    """
    try:
        loop = asyncio.get_event_loop()
        
        async def read_text():
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(1000)
        
        text = await loop.run_in_executor(None, read_text)
        return _create_text_preview(text, "Text File", max_width, max_height)
    except Exception as e:
        logger.error(f"Failed to read text file: {e}")
        return _create_error_preview(str(e), max_width, max_height)


async def _resize_and_encode_image(img: Image.Image, max_width: int, max_height: int) -> bytes:
    """
    Resizes image to fit within max dimensions and encodes as JPEG.
    """
    # Resize maintaining aspect ratio
    img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
    
    # Encode to JPEG
    output = io.BytesIO()
    img.save(output, format='JPEG', quality=85)
    output.seek(0)
    return output.getvalue()


def _create_text_preview(text: str, title: str, width: int, height: int) -> bytes:
    """
    Creates an image with text content for preview.
    """
    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    # Try to use a better font, fallback to default
    try:
        from PIL import ImageFont
        font = ImageFont.load_default()
    except:
        font = ImageFont.load_default()
    
    # Draw title
    title_text = f"📄 {title}"
    draw.text((10, 10), title_text, fill=(0, 0, 0), font=font)
    
    # Draw preview text
    y_position = 40
    lines = text.split('\n')
    for line in lines[:30]:  # Limit to 30 lines
        if y_position > height - 20:
            draw.text((10, y_position), "...", fill=(128, 128, 128), font=font)
            break
        draw.text((10, y_position), line[:60], fill=(50, 50, 50), font=font)
        y_position += 15
    
    # Encode
    output = io.BytesIO()
    img.save(output, format='JPEG', quality=85)
    output.seek(0)
    return output.getvalue()


def _create_file_info_preview(file_type: str, filepath: str, width: int, height: int) -> bytes:
    """
    Creates a simple info preview for unsupported file types.
    """
    img = Image.new('RGB', (width, height), color=(240, 240, 240))
    draw = ImageDraw.Draw(img)
    
    try:
        from PIL import ImageFont
        font = ImageFont.load_default()
    except:
        font = ImageFont.load_default()
    
    file_size = os.path.getsize(filepath)
    file_name = os.path.basename(filepath)
    
    text_lines = [
        f"📋 {file_type}",
        f"Файл: {file_name}",
        f"Размер: {file_size / 1024:.1f} KB",
        "",
        "Файл готов к печати.",
        "Нажми ✅ Печать, чтобы начать."
    ]
    
    y = 50
    for line in text_lines:
        draw.text((30, y), line, fill=(0, 0, 0), font=font)
        y += 40
    
    output = io.BytesIO()
    img.save(output, format='JPEG', quality=85)
    output.seek(0)
    return output.getvalue()


def _create_unsupported_preview(file_ext: str, width: int, height: int) -> bytes:
    """
    Creates preview for unsupported file type.
    """
    img = Image.new('RGB', (width, height), color=(255, 200, 200))
    draw = ImageDraw.Draw(img)
    
    try:
        from PIL import ImageFont
        font = ImageFont.load_default()
    except:
        font = ImageFont.load_default()
    
    draw.text((50, 150), "❌ Формат не поддерживается", fill=(200, 0, 0), font=font)
    draw.text((100, 250), file_ext.upper(), fill=(100, 0, 0), font=font)
    
    output = io.BytesIO()
    img.save(output, format='JPEG', quality=85)
    output.seek(0)
    return output.getvalue()


def _create_error_preview(error_msg: str, width: int, height: int) -> bytes:
    """
    Creates preview showing an error.
    """
    img = Image.new('RGB', (width, height), color=(255, 100, 100))
    draw = ImageDraw.Draw(img)
    
    try:
        from PIL import ImageFont
        font = ImageFont.load_default()
    except:
        font = ImageFont.load_default()
    
    draw.text((30, 150), "⚠️ Ошибка при создании", fill=(200, 0, 0), font=font)
    draw.text((30, 250), "превью документа", fill=(200, 0, 0), font=font)
    
    output = io.BytesIO()
    img.save(output, format='JPEG', quality=85)
    output.seek(0)
    return output.getvalue()


def apply_print_settings(image: Image.Image, settings: dict) -> Image.Image:
    """
    Applies print settings to an image for preview.
    """
    width, height = image.size
    
    # Handle orientation
    orientation = settings.get("orientation", "auto")
    if orientation == "auto":
        # Auto-detect best orientation based on image aspect ratio
        if width > height:
            orientation = "landscape"
        else:
            orientation = "portrait"
    
    # For landscape images, auto-rotate if needed
    if orientation == "landscape" and width < height:
        image = image.rotate(90, expand=True)
    elif orientation == "portrait" and width > height:
        image = image.rotate(90, expand=True)
    
    # Apply color mode
    color_mode = settings.get("color_mode", "color")
    if color_mode == "grayscale":
        image = image.convert('L').convert('RGB')
    elif color_mode == "black_white":
        image = image.convert('L')
        image = Image.new('RGB', image.size, (255, 255, 255))
        pixels = image.load()
        bw_pixels = Image.new('L', image.size).load()
        for i in range(image.size[0]):
            for j in range(image.size[1]):
                bw_pixels[i, j] = 200 if i % 2 == j % 2 else 50
    
    # Apply scaling
    scale = settings.get("scale", "fit")
    custom_scale = settings.get("custom_scale", 100)
    
    if scale == "custom":
        new_width = int(image.width * custom_scale / 100)
        new_height = int(image.height * custom_scale / 100)
        image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
    # fit and fill are applied during printing, not preview
    
    return image
